import sys
import subprocess

# Garante que o python-pptx está instalado
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt

def criar_apresentacao():
    prs = Presentation()
    img_path = "images/"  # ajuste para o diretório de ícones/imagens

    slides_info = [
        # (layout_index, título, subtítulo, nome_imagem, [bullets])
        (0, "[Nome da Sua Agência]", "Consultoria de Comunicação e Conteúdo Digital", "logo.png", []),
        (1, "O Problema que Resolvemos", None, "problema.png", [
            "Pequenas empresas lutam para se destacar online",
            "Falta de tempo e conhecimento para criar conteúdo estratégico",
            "Dificuldade em transformar presença digital em resultados concretos",
            "Esforços dispersos sem um plano claro",
        ]),
        (1, "Nossa Solução", None, "solucao.png", [
            "Conteúdo estratégico para redes sociais e canais digitais",
            "Parceria próxima e personalizada",
            "Transformação da presença digital em crescimento real",
        ]),
        (1, "Visão Geral de Serviços", None, "servicos.png", [
            "Consultoria e Estratégia Digital",
            "Produção de Conteúdo e Gestão de Canais",
            "Performance e Análise de Resultados",
            "Serviços Complementares",
        ]),
        (1, "Consultoria e Estratégia Digital", "Criando as Bases do Sucesso Digital", "estrategia.png", [
            "Diagnóstico e Análise de Marca, Mercado e Público",
            "Definição de Persona e Jornada do Cliente",
            "Planejamento Estratégico de Conteúdo",
            "Definição de Tom de Voz",
            "Consultoria contínua de comunicação",
        ]),
        (1, "Produção & Gestão de Conteúdo", "Sua Marca Ganhando Voz e Visibilidade", "producao.png", [
            "Posts para Instagram, LinkedIn e outras redes",
            "Artigos e reportagens para blog/portal",
            "Campanhas de e-mail marketing",
            "Design de conteúdo visual",
            "Gestão editorial e revisão",
        ]),
        (1, "Performance & Análise", "Transformando Conteúdo em Crescimento", "performance.png", [
            "Gestão de tráfego pago (ads)",
            "Monitoramento de menções e performance",
            "Relatórios de resultados e otimizações contínuas",
        ]),
        (1, "Serviços Complementares", "Agregando Valor à Comunicação", "complementares.png", [
            "Assessoria de imprensa digital",
            "Cobertura de eventos em tempo real",
        ]),
        (1, "Por Que Nos Escolher?", None, "porque.png", [
            "13 anos de experiência em comunicação estratégica",
            "Foco em resultados concretos",
            "Parceria próxima com pequenas empresas",
            "Visão completa: da estratégia à performance",
        ]),
        (1, "Vamos Conversar?", None, "contato.png", [
            "Telefone: [Seu Telefone]",
            "E-mail: [Seu E-mail]",
            "Site/Portfólio: [URL]",
            "LinkedIn/Instagram: @[perfil]",
        ]),
    ]

    for layout_idx, title, subtitle, image_file, bullets in slides_info:
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        slide.shapes.title.text = title

        # subtítulo (quando existir)
        if subtitle:
            try:
                slide.placeholders[1].text = subtitle
            except Exception:
                pass

        # bullets
        if bullets:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for i, item in enumerate(bullets):
                p = tf.add_paragraph() if i else tf.paragraphs[0]
                p.text = item
                p.level = 0
                p.font.size = Pt(18)

        # imagem/ícone (ajuste posição conforme seu template)
        try:
            slide.shapes.add_picture(img_path + image_file, Inches(6), Inches(4.5), width=Inches(2))
        except FileNotFoundError:
            pass  # se faltar a imagem, ignora

    prs.save("apresentacao_agencia.pptx")
    print("Apresentação gerada: apresentacao_agencia.pptx")

# if _name_ == "_main_":
criar_apresentacao()